"""
图片提取工具 - Word文档服务器

这个模块提供了从Word/PDF/PPT/Excel文件中提取图片并打包为ZIP文件的功能。

支持的文件格式：
- Word文档 (.docx, .doc)
- PDF文件 (.pdf)
- PowerPoint演示文稿 (.pptx, .ppt)
- Excel电子表格 (.xlsx, .xls)

功能特点：
1. 支持多种文件格式的图片提取
2. 自动下载远程文件URL
3. 提取所有图片并保存为常见格式
4. 自动打包为ZIP文件
5. 提供详细的提取统计信息
6. 完善的错误处理机制

使用方式：
1. 直接调用 extract_images_from_file() 提取图片
2. 调用 extract_images_and_upload() 提取并上传ZIP文件
3. 使用 ImageExtractor 类进行更精细的控制
"""

import os
import zipfile
import tempfile
import shutil
import uuid
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import requests
from PIL import Image
import io

# 导入项目内的工具函数
from word_document_server.utils.file_utils import download_file_from_url, upload_file_to_server


class ImageExtractor:
    """
    图片提取器 - 核心类
    
    这个类实现了从各种文档格式中提取图片的功能，支持Word、PDF、PPT、Excel等格式。
    
    工作原理：
    1. 根据文件扩展名识别文件类型
    2. 使用相应的库解析文档内容
    3. 提取所有嵌入的图片
    4. 将图片保存为常见格式（PNG、JPG等）
    5. 打包为ZIP文件
    
    支持的文件格式：
    - Word文档：.docx, .doc
    - PDF文件：.pdf
    - PowerPoint：.pptx, .ppt
    - Excel：.xlsx, .xls
    
    图片格式支持：
    - 输入：各种嵌入格式（PNG、JPG、GIF、BMP等）
    - 输出：PNG、JPG格式
    """
    
    def __init__(self):
        """初始化图片提取器"""
        self.extracted_images = []
        self.temp_dir = None
        self.stats = {
            "total_images": 0,
            "successful_extractions": 0,
            "failed_extractions": 0,
            "errors": []
        }
    
    def _create_temp_directory(self) -> str:
        """创建临时目录用于存储提取的图片"""
        if self.temp_dir is None:
            self.temp_dir = tempfile.mkdtemp(prefix="image_extraction_")
        return self.temp_dir
    
    def _get_file_extension(self, file_path: str) -> str:
        """获取文件扩展名（小写）"""
        return Path(file_path).suffix.lower()
    
    def _is_supported_format(self, file_path: str) -> bool:
        """检查文件格式是否支持"""
        ext = self._get_file_extension(file_path)
        supported_formats = ['.docx', '.doc', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls']
        return ext in supported_formats
    
    def _extract_from_docx(self, file_path: str) -> List[str]:
        """从Word文档(.docx)中提取图片"""
        try:
            from docx import Document
            doc = Document(file_path)
            extracted_paths = []
            
            # 提取文档中的图片
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        image_name = f"image_{len(extracted_paths) + 1}.png"
                        image_path = os.path.join(self._create_temp_directory(), image_name)
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        extracted_paths.append(image_path)
                        self.stats["successful_extractions"] += 1
                        
                    except Exception as e:
                        error_msg = f"提取Word图片失败: {str(e)}"
                        self.stats["errors"].append(error_msg)
                        self.stats["failed_extractions"] += 1
            
            return extracted_paths
            
        except Exception as e:
            error_msg = f"处理Word文档失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def _extract_from_pdf(self, file_path: str) -> List[str]:
        """从PDF文件中提取图片"""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            extracted_paths = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # 灰度或RGB
                            image_data = pix.tobytes("png")
                        else:  # CMYK或其他格式
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            image_data = pix1.tobytes("png")
                            pix1 = None
                        
                        image_name = f"page_{page_num + 1}_image_{img_index + 1}.png"
                        image_path = os.path.join(self._create_temp_directory(), image_name)
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        extracted_paths.append(image_path)
                        self.stats["successful_extractions"] += 1
                        
                        pix = None
                        
                    except Exception as e:
                        error_msg = f"提取PDF图片失败 (页面{page_num + 1}, 图片{img_index + 1}): {str(e)}"
                        self.stats["errors"].append(error_msg)
                        self.stats["failed_extractions"] += 1
            
            doc.close()
            return extracted_paths
            
        except ImportError:
            error_msg = "PyMuPDF库未安装，无法处理PDF文件。请安装: pip install PyMuPDF"
            self.stats["errors"].append(error_msg)
            return []
        except Exception as e:
            error_msg = f"处理PDF文件失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def _extract_from_pptx(self, file_path: str) -> List[str]:
        """从PowerPoint演示文稿(.pptx)中提取图片"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            extracted_paths = []
            
            for slide_num, slide in enumerate(prs.slides):
                for shape_num, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'image'):
                        try:
                            image_data = shape.image.blob
                            image_name = f"slide_{slide_num + 1}_image_{shape_num + 1}.png"
                            image_path = os.path.join(self._create_temp_directory(), image_name)
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_data)
                            
                            extracted_paths.append(image_path)
                            self.stats["successful_extractions"] += 1
                            
                        except Exception as e:
                            error_msg = f"提取PPT图片失败 (幻灯片{slide_num + 1}, 形状{shape_num + 1}): {str(e)}"
                            self.stats["errors"].append(error_msg)
                            self.stats["failed_extractions"] += 1
            
            return extracted_paths
            
        except ImportError:
            error_msg = "python-pptx库未安装，无法处理PPT文件。请安装: pip install python-pptx"
            self.stats["errors"].append(error_msg)
            return []
        except Exception as e:
            error_msg = f"处理PPT文件失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def _extract_from_xlsx(self, file_path: str) -> List[str]:
        """从Excel电子表格(.xlsx)中提取图片"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path)
            extracted_paths = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                if hasattr(ws, '_images'):
                    for img_index, img in enumerate(ws._images):
                        try:
                            image_data = img._data()
                            image_name = f"sheet_{sheet_name}_image_{img_index + 1}.png"
                            image_path = os.path.join(self._create_temp_directory(), image_name)
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_data)
                            
                            extracted_paths.append(image_path)
                            self.stats["successful_extractions"] += 1
                            
                        except Exception as e:
                            error_msg = f"提取Excel图片失败 (工作表{sheet_name}, 图片{img_index + 1}): {str(e)}"
                            self.stats["errors"].append(error_msg)
                            self.stats["failed_extractions"] += 1
            
            return extracted_paths
            
        except ImportError:
            error_msg = "openpyxl库未安装，无法处理Excel文件。请安装: pip install openpyxl"
            self.stats["errors"].append(error_msg)
            return []
        except Exception as e:
            error_msg = f"处理Excel文件失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def _extract_from_legacy_formats(self, file_path: str) -> List[str]:
        """从旧格式文件(.doc, .ppt, .xls)中提取图片"""
        # 对于旧格式，我们尝试使用一些通用方法
        try:
            # 尝试使用zipfile读取（某些旧格式实际上是ZIP容器）
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                extracted_paths = []
                for file_info in zip_ref.filelist:
                    if any(ext in file_info.filename.lower() for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']):
                        try:
                            image_data = zip_ref.read(file_info.filename)
                            image_name = f"legacy_{len(extracted_paths) + 1}{Path(file_info.filename).suffix}"
                            image_path = os.path.join(self._create_temp_directory(), image_name)
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_data)
                            
                            extracted_paths.append(image_path)
                            self.stats["successful_extractions"] += 1
                            
                        except Exception as e:
                            error_msg = f"提取旧格式图片失败: {str(e)}"
                            self.stats["errors"].append(error_msg)
                            self.stats["failed_extractions"] += 1
                
                return extracted_paths
                
        except zipfile.BadZipFile:
            error_msg = "不支持的文件格式或文件损坏"
            self.stats["errors"].append(error_msg)
            return []
        except Exception as e:
            error_msg = f"处理旧格式文件失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def extract_images(self, file_path: str) -> List[str]:
        """
        从文件中提取所有图片
        
        Args:
            file_path: 文件路径
            
        Returns:
            List[str]: 提取的图片文件路径列表
        """
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            self.stats["errors"].append(error_msg)
            return []
        
        if not self._is_supported_format(file_path):
            error_msg = f"不支持的文件格式: {self._get_file_extension(file_path)}"
            self.stats["errors"].append(error_msg)
            return []
        
        # 重置统计信息
        self.stats = {
            "total_images": 0,
            "successful_extractions": 0,
            "failed_extractions": 0,
            "errors": []
        }
        
        ext = self._get_file_extension(file_path)
        extracted_paths = []
        
        try:
            if ext == '.docx':
                extracted_paths = self._extract_from_docx(file_path)
            elif ext == '.pdf':
                extracted_paths = self._extract_from_pdf(file_path)
            elif ext == '.pptx':
                extracted_paths = self._extract_from_pptx(file_path)
            elif ext == '.xlsx':
                extracted_paths = self._extract_from_xlsx(file_path)
            elif ext in ['.doc', '.ppt', '.xls']:
                extracted_paths = self._extract_from_legacy_formats(file_path)
            
            self.stats["total_images"] = len(extracted_paths)
            self.extracted_images = extracted_paths
            
            return extracted_paths
            
        except Exception as e:
            error_msg = f"提取图片过程中发生错误: {str(e)}"
            self.stats["errors"].append(error_msg)
            return []
    
    def create_zip_archive(self, output_filename: str = None) -> str:
        """
        将提取的图片打包为ZIP文件
        
        Args:
            output_filename: 输出ZIP文件名，如果为None则自动生成
            
        Returns:
            str: ZIP文件路径
        """
        if not self.extracted_images:
            raise ValueError("没有图片可以打包")
        
        if output_filename is None:
            output_filename = f"extracted_images_{uuid.uuid4().hex[:8]}.zip"
        
        if not output_filename.endswith('.zip'):
            output_filename += '.zip'
        
        try:
            with zipfile.ZipFile(output_filename, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for image_path in self.extracted_images:
                    if os.path.exists(image_path):
                        # 使用相对路径作为ZIP内的文件名
                        arcname = os.path.basename(image_path)
                        zipf.write(image_path, arcname)
            
            return output_filename
            
        except Exception as e:
            error_msg = f"创建ZIP文件失败: {str(e)}"
            self.stats["errors"].append(error_msg)
            raise
    
    def cleanup(self):
        """清理临时文件和目录"""
        if self.temp_dir and os.path.exists(self.temp_dir):
            try:
                shutil.rmtree(self.temp_dir)
                self.temp_dir = None
            except Exception as e:
                logging.warning(f"清理临时目录失败: {str(e)}")


def extract_images_from_file(
    file_path: str = None,
    file_url: str = None,
    output_filename: str = None
) -> Dict[str, Any]:
    """
    从文件中提取图片并打包为ZIP文件
    
    Args:
        file_path: 本地文件路径（与file_url二选一）
        file_url: 远程文件URL（与file_path二选一）
        output_filename: 输出ZIP文件名（可选，自动生成）
        
    Returns:
        Dict[str, Any]: 包含以下字段的结果字典：
            - success: 是否成功
            - message: 处理结果描述
            - zip_file_path: ZIP文件路径
            - extracted_images: 提取的图片列表
            - stats: 统计信息
            - error: 如果出错，包含错误信息
    """
    extractor = ImageExtractor()
    local_file_path = None
    
    try:
        # 处理文件路径或URL
        if file_url:
            if not file_path:
                # 从URL下载文件
                temp_dir = tempfile.mkdtemp(prefix="download_")
                local_file_path = download_file_from_url(file_url, temp_dir)
            else:
                local_file_path = file_path
        elif file_path:
            local_file_path = file_path
        else:
            return {
                "success": False,
                "error": "必须提供file_path或file_url参数"
            }
        
        # 提取图片
        extracted_images = extractor.extract_images(local_file_path)
        
        if not extracted_images:
            return {
                "success": False,
                "error": "未找到任何图片",
                "stats": extractor.stats
            }
        
        # 创建ZIP文件
        zip_file_path = extractor.create_zip_archive(output_filename)
        
        return {
            "success": True,
            "message": f"成功提取{len(extracted_images)}张图片并打包为ZIP文件",
            "zip_file_path": zip_file_path,
            "extracted_images": extracted_images,
            "stats": extractor.stats
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"提取图片失败: {str(e)}",
            "stats": extractor.stats if hasattr(extractor, 'stats') else {}
        }
    finally:
        extractor.cleanup()


def extract_images_and_upload(
    file_path: str = None,
    file_url: str = None,
    output_filename: str = None
) -> Dict[str, Any]:
    """
    从文件中提取图片，打包为ZIP文件并上传到服务器
    
    Args:
        file_path: 本地文件路径（与file_url二选一）
        file_url: 远程文件URL（与file_path二选一）
        output_filename: 输出ZIP文件名（可选，自动生成）
        
    Returns:
        Dict[str, Any]: 包含以下字段的结果字典：
            - success: 是否成功
            - message: 处理结果描述
            - public_url: 公网下载链接
            - remote_path: 服务器上的文件路径
            - zip_file_path: 本地ZIP文件路径
            - stats: 统计信息
            - error: 如果出错，包含错误信息
    """
    try:
        # 首先提取图片并创建ZIP文件
        extract_result = extract_images_from_file(file_path, file_url, output_filename)
        
        if not extract_result["success"]:
            return extract_result
        
        zip_file_path = extract_result["zip_file_path"]
        
        # 上传到服务器
        REMOTE_DIR = "/root/files"
        SERVER = "8.156.74.79"
        USERNAME = "root"
        PASSWORD = "zfsZBC123."
        remote_path = os.path.join(REMOTE_DIR, os.path.basename(zip_file_path))
        
        upload_result = upload_file_to_server(zip_file_path, remote_path, SERVER, USERNAME, PASSWORD)
        public_url = f"http://8.156.74.79:8001/{os.path.basename(zip_file_path)}"
        
        return {
            "success": True,
            "message": "图片提取并上传成功",
            "public_url": public_url,
            "remote_path": remote_path,
            "zip_file_path": zip_file_path,
            "upload_result": upload_result,
            "stats": extract_result["stats"]
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"上传失败: {str(e)}",
            "stats": extract_result.get("stats", {}) if 'extract_result' in locals() else {}
        }


def get_supported_formats() -> Dict[str, Any]:
    """
    获取支持的文件格式信息
    
    Returns:
        Dict[str, Any]: 支持的文件格式信息
    """
    return {
        "supported_formats": [
            {
                "extension": ".docx",
                "description": "Word文档（新格式）",
                "features": ["文本", "图片", "表格", "格式"]
            },
            {
                "extension": ".doc",
                "description": "Word文档（旧格式）",
                "features": ["文本", "图片", "表格", "格式"]
            },
            {
                "extension": ".pdf",
                "description": "PDF文档",
                "features": ["文本", "图片", "矢量图形", "表单"]
            },
            {
                "extension": ".pptx",
                "description": "PowerPoint演示文稿（新格式）",
                "features": ["幻灯片", "图片", "动画", "图表"]
            },
            {
                "extension": ".ppt",
                "description": "PowerPoint演示文稿（旧格式）",
                "features": ["幻灯片", "图片", "动画", "图表"]
            },
            {
                "extension": ".xlsx",
                "description": "Excel电子表格（新格式）",
                "features": ["数据", "图表", "图片", "公式"]
            },
            {
                "extension": ".xls",
                "description": "Excel电子表格（旧格式）",
                "features": ["数据", "图表", "图片", "公式"]
            }
        ],
        "image_formats": ["PNG", "JPG", "JPEG", "GIF", "BMP", "TIFF"],
        "output_format": "ZIP压缩包"
    }
